import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Set slide size to widescreen (16:9)
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def generate_lecture_ppt(self, book_data, lecture_plan_json):
        """
        책 정보와 강의안을 바탕으로 PPT를 생성합니다.
        """
        try:
            lecture_plan = json.loads(lecture_plan_json) if isinstance(lecture_plan_json, str) else lecture_plan_json
            
            # Title slide
            self._create_title_slide(book_data, lecture_plan)
            
            # Overview slide
            if lecture_plan.get('lecture_overview'):
                self._create_overview_slide(lecture_plan['lecture_overview'])
            
            # Individual lecture slides
            if lecture_plan.get('lectures'):
                for lecture in lecture_plan['lectures']:
                    self._create_lecture_slide(lecture)
            
            # Assessment and resources slide
            if lecture_plan.get('assessment') or lecture_plan.get('resources'):
                self._create_assessment_resources_slide(lecture_plan.get('assessment'), lecture_plan.get('resources'))
            
            return True
            
        except Exception as e:
            print(f"PPT 생성 중 오류 발생: {str(e)}")
            return False
    
    def _create_title_slide(self, book_data, lecture_plan):
        """제목 슬라이드 생성"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Set title
        if lecture_plan.get('lecture_overview', {}).get('title'):
            title.text = lecture_plan['lecture_overview']['title']
        else:
            title.text = f"{book_data.get('title', '')} 강의안"
        
        # Set subtitle with book info
        subtitle_text = []
        if book_data.get('author'):
            subtitle_text.append(f"저자: {book_data['author']}")
        if lecture_plan.get('lecture_overview', {}).get('target_audience'):
            subtitle_text.append(f"대상: {lecture_plan['lecture_overview']['target_audience']}")
        if lecture_plan.get('lecture_overview', {}).get('duration'):
            subtitle_text.append(f"총 시간: {lecture_plan['lecture_overview']['duration']}")
        
        subtitle.text = "\n".join(subtitle_text)
        
        # Style title
        self._style_title(title)
        self._style_subtitle(subtitle)
    
    def _create_overview_slide(self, overview):
        """강의 개요 슬라이드 생성"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "강의 개요"
        self._style_title(title)
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Description
        if overview.get('description'):
            p = tf.paragraphs[0]
            p.text = overview['description']
            p.font.size = Pt(18)
            
        # Learning objectives
        if overview.get('learning_objectives'):
            tf.add_paragraph()
            p = tf.add_paragraph()
            p.text = "학습 목표:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for objective in overview['learning_objectives']:
                p = tf.add_paragraph()
                p.text = f"• {objective}"
                p.font.size = Pt(14)
                p.level = 1
    
    def _create_lecture_slide(self, lecture):
        """개별 강의 슬라이드 생성"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{lecture.get('lecture_number', '')}강: {lecture.get('title', '')}"
        self._style_title(title)
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Duration and objectives
        info_text = []
        if lecture.get('duration'):
            info_text.append(f"수업 시간: {lecture['duration']}")
        
        if info_text:
            p = tf.paragraphs[0]
            p.text = " | ".join(info_text)
            p.font.size = Pt(14)
            p.font.italic = True
            tf.add_paragraph()
        
        # Objectives
        if lecture.get('objectives'):
            p = tf.add_paragraph()
            p.text = "학습 목표:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for obj in lecture['objectives']:
                p = tf.add_paragraph()
                p.text = f"• {obj}"
                p.font.size = Pt(14)
                p.level = 1
            
            tf.add_paragraph()
        
        # Outline
        if lecture.get('outline'):
            p = tf.add_paragraph()
            p.text = "강의 구성:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for section in lecture['outline']:
                p = tf.add_paragraph()
                section_text = f"• {section.get('section', '')}"
                if section.get('time'):
                    section_text += f" ({section['time']})"
                p.text = section_text
                p.font.size = Pt(14)
                p.level = 1
                
                if section.get('content'):
                    p = tf.add_paragraph()
                    p.text = f"  - {section['content']}"
                    p.font.size = Pt(12)
                    p.level = 2
        
        # Key concepts
        if lecture.get('key_concepts'):
            tf.add_paragraph()
            p = tf.add_paragraph()
            p.text = "핵심 개념:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            p = tf.add_paragraph()
            p.text = " | ".join(lecture['key_concepts'])
            p.font.size = Pt(14)
            p.level = 1
    
    def _create_assessment_resources_slide(self, assessment, resources):
        """평가 및 자료 슬라이드 생성"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "평가 방법 및 학습 자료"
        self._style_title(title)
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Assessment
        if assessment:
            p = tf.paragraphs[0]
            p.text = "평가 방법:"
            p.font.size = Pt(18)
            p.font.bold = True
            
            if assessment.get('methods'):
                for method in assessment['methods']:
                    p = tf.add_paragraph()
                    p.text = f"• {method}"
                    p.font.size = Pt(14)
                    p.level = 1
            
            if assessment.get('criteria'):
                p = tf.add_paragraph()
                p.text = f"평가 기준: {assessment['criteria']}"
                p.font.size = Pt(12)
                p.font.italic = True
            
            tf.add_paragraph()
        
        # Resources
        if resources:
            p = tf.add_paragraph()
            p.text = "학습 자료:"
            p.font.size = Pt(18)
            p.font.bold = True
            
            if resources.get('required'):
                p = tf.add_paragraph()
                p.text = "필수 자료:"
                p.font.size = Pt(16)
                p.font.bold = True
                
                for resource in resources['required']:
                    p = tf.add_paragraph()
                    p.text = f"• {resource}"
                    p.font.size = Pt(14)
                    p.level = 1
            
            if resources.get('recommended'):
                p = tf.add_paragraph()
                p.text = "추천 자료:"
                p.font.size = Pt(16)
                p.font.bold = True
                
                for resource in resources['recommended']:
                    p = tf.add_paragraph()
                    p.text = f"• {resource}"
                    p.font.size = Pt(14)
                    p.level = 1
    
    def _style_title(self, title_shape):
        """제목 스타일 설정"""
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _style_subtitle(self, subtitle_shape):
        """부제목 스타일 설정"""
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.alignment = PP_ALIGN.CENTER
    
    def save_ppt(self, filename):
        """PPT 파일 저장"""
        try:
            # Ensure directory exists
            os.makedirs(os.path.dirname(filename), exist_ok=True)
            self.prs.save(filename)
            return True
        except Exception as e:
            print(f"PPT 저장 중 오류 발생: {str(e)}")
            return False